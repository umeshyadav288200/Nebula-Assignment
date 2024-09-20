import fitz  # PyMuPDF for PDF processing
import pandas as pd
from pptx import Presentation

# Extract text from PDF
def extract_text_from_pdf(uploaded_file):
    document_text = ""
    pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        document_text += page.get_text("text")
    pdf_document.close()
    return document_text

# Extract text from spreadsheets (Excel or CSV)
def extract_text_from_spreadsheet(uploaded_file):
    if uploaded_file.name.endswith(".csv"):
        df = pd.read_csv(uploaded_file)
    else:
        df = pd.read_excel(uploaded_file)
    return df.to_string()

# Extract text from PowerPoint
def extract_text_from_ppt(uploaded_file):
    document_text = ""
    prs = Presentation(uploaded_file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                document_text += shape.text + "\n"
    return document_text
